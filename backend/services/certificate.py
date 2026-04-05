"""
Certificate generation service module.
Uses python-pptx to edit PPTX template and converts to PDF.
"""
import os
import uuid
import subprocess
from datetime import datetime
from config import CERTIFICATES_DIRECTORY, logger

# Check for required libraries
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    PYPPTX_AVAILABLE = True
except ImportError:
    PYPPTX_AVAILABLE = False
    logger.warning("python-pptx not installed. Run: pip install python-pptx")

# Path to LibreOffice for conversion (fallback if comtypes not available)
LIBREOFFICE_PATH = "soffice"  # Usually in PATH on Linux, or specific path on Windows


def _replace_text_in_shape(shape, replacements):
    """
    Replace multiple placeholders in a shape while preserving font formatting.
    replacements: dict of {old_text: new_text}
    """
    if not hasattr(shape, "text_frame"):
        return False
    
    # Get all text from shape
    full_text = shape.text
    if not any(placeholder in full_text for placeholder in replacements.keys()):
        return False
    
    # Process each paragraph
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            for old_text, new_text in replacements.items():
                if old_text in run.text:
                    # Get font properties before replacement
                    font_name = run.font.name
                    font_size = run.font.size
                    font_bold = run.font.bold
                    font_italic = run.font.italic
                    font_color = None
                    try:
                        if run.font.color and run.font.color.rgb:
                            font_color = run.font.color.rgb
                    except:
                        pass
                    
                    # Replace the text
                    run.text = run.text.replace(old_text, new_text)
                    
                    # Re-apply font properties (they might reset)
                    try:
                        if font_name:
                            run.font.name = font_name
                        if font_size:
                            run.font.size = font_size
                        if font_bold is not None:
                            run.font.bold = font_bold
                        if font_italic is not None:
                            run.font.italic = font_italic
                        if font_color:
                            run.font.color.rgb = font_color
                    except:
                        pass
                    break  # Move to next run after successful replacement
    return True


def _get_output_directory():
    """Get the output directory for generated certificates."""
    # Use backend directory as base (parent of services/)
    backend_dir = os.path.dirname(os.path.dirname(__file__))
    output_dir = os.path.join(backend_dir, "static", "certificates", "generate")
    abs_output_dir = os.path.abspath(output_dir)
    return abs_output_dir


def _get_title_prefix(project_title: str) -> str:
    """
    Generate certificate ID prefix based on project title.
    - "Trial Bootcamp Data Science & AI" → TBDSAI
    - "Trial Bootcamp Data Science" → TBDS
    - "Trial Bootcamp Artificial Intelligence" or "Trial Bootcamp AI" → TBAI
    """
    title_upper = project_title.upper()
    
    if "DATA SCIENCE & AI" in title_upper or "DATA SCIENCE AND AI" in title_upper:
        return "TBDSAI"
    elif "DATA SCIENCE" in title_upper:
        return "TBDS"
    elif "ARTIFICIAL INTELLIGENCE" in title_upper or ("TRIAL BOOTCAMP AI" in title_upper and "DATA SCIENCE" not in title_upper):
        return "TBAI"
    else:
        return "TBDSAI"  # Default


def _parse_date_for_id(date_str: str) -> str:
    """
    Parse execution date string to get MMYY format.
    Supports: "October 2025", "23 - 27 October 2025", "10/2025", etc.
    Returns: "MMYY" format
    """
    import re
    
    date_str = date_str.strip()
    
    # Try to extract month and year from various formats
    # Format: "23 - 27 October 2025" or "23-27 October 2025"
    match = re.search(r'([A-Za-z]+)\s+(\d{4})', date_str)
    if match:
        month_name = match.group(1)
        year = match.group(2)
        date_str = f"{month_name} {year}"
    
    # Try to parse common date formats
    date_formats = ["%B %Y", "%b %Y", "%m/%Y", "%Y-%m", "%Y"]
    
    for fmt in date_formats:
        try:
            parsed = datetime.strptime(date_str.strip(), fmt)
            return parsed.strftime("%m%y")
        except ValueError:
            continue
    
    # If parsing fails, use current date
    return datetime.now().strftime("%m%y")


def _format_date_range(start_date: str, end_date: str) -> str:
    """
    Format date range intelligently.
    - Same month: "23 to 27 October 2025"
    - Different months: "23 October to 27 November 2025"
    
    Accepts both ISO format (2026-03-01) and text format (23 October 2025)
    """
    from datetime import datetime
    import re
    
    def parse_date(date_str: str) -> tuple:
        """Parse date string to (day, month_name, year)"""
        date_str = date_str.strip()
        
        # Try ISO format: "2026-03-01"
        try:
            date_obj = datetime.strptime(date_str, "%Y-%m-%d")
            day = str(date_obj.day)
            month_name = date_obj.strftime("%B")
            year = str(date_obj.year)
            return (day, month_name, year)
        except ValueError:
            pass
        
        # Try format: "23 October 2025"
        match = re.match(r'(\d{1,2})\s+([A-Za-z]+)\s+(\d{4})', date_str)
        if match:
            return match.groups()
        
        # Try format: "October 2025" (month only)
        match = re.match(r'([A-Za-z]+)\s+(\d{4})', date_str)
        if match:
            return (None, match.group(1), match.group(2))
        
        return None, None, None
    
    start_day, start_month, start_year = parse_date(start_date)
    end_day, end_month, end_year = parse_date(end_date)
    
    if start_day is None or end_day is None:
        # Fallback: just use simple format
        return f"{start_date} to {end_date}"
    
    # Check if same year and same month
    if start_year == end_year and start_month == end_month:
        return f"{start_day} to {end_day} {start_month} {start_year}"
    else:
        return f"{start_day} {start_month} to {end_day} {end_month} {end_year}"


def _parse_certificate_info(certificate_info: str) -> tuple:
    """
    Parse combined certificate info string.
    Format: "Trial Bootcamp Data Science & AI - Intelligo ID - 23 - 27 October 2025"
    Returns: (project_title, execution_date)
    """
    import re
    
    # Pattern to detect if the last part is a date range like "23 - 27 October 2025"
    # or a single date like "October 2025"
    date_pattern = re.compile(r'^(?:\d+\s*-\s*)?\d+\s+[A-Za-z]+\s+\d{4}$')
    
    # Split by " - " but be careful with date ranges
    # First, try to find if there's a date at the end
    parts = certificate_info.split(" - ")
    
    if len(parts) >= 2:
        # Check if last part looks like a date
        last_part = parts[-1].strip()
        
        if date_pattern.match(last_part):
            # Last part is a date (single or range)
            execution_date = last_part
            project_title = " - ".join(parts[:-1]).strip()
        else:
            # Could be "23 - 27 October 2025" where parts[-2] is part of date
            if len(parts) >= 3:
                # Check if last two parts form a date range
                second_last = parts[-2].strip()
                # If second_last is just a number (like "23"), it's part of date range
                if second_last.isdigit() and len(second_last) <= 2:
                    # Reconstruct date range
                    execution_date = f"{parts[-2]} - {parts[-1]}"
                    project_title = " - ".join(parts[:-2]).strip()
                else:
                    # Just use last part as date
                    execution_date = last_part
                    project_title = " - ".join(parts[:-1]).strip()
            else:
                execution_date = last_part
                project_title = " - ".join(parts[:-1]).strip()
    else:
        # Single part, use as-is
        project_title = certificate_info
        execution_date = datetime.now().strftime("%B %Y")
    
    return project_title, execution_date


def generate_certificate(name: str, program_title: str, start_date: str, end_date: str) -> str:
    """
    Generate a PDF certificate by editing the PPTX template and converting to PDF.
    
    Args:
        name: Participant name
        program_title: Program name (e.g., "Trial Bootcamp Data Science & AI - Intelligo ID")
        start_date: Start date (e.g., "23 October 2025")
        end_date: End date (e.g., "27 October 2025")
    
    Dynamic placeholders replaced:
    - {{NAMA}} - Participant name
    - {{JUDUL}} - Program title
    - {{PELAKSANAAN}} - Date range formatted as "23 to 27 October 2025"
    - {{id}} - Certificate ID
    """
    if not PYPPTX_AVAILABLE:
        raise RuntimeError("python-pptx is required. Install with: pip install python-pptx")
    
    # Format execution date intelligently
    execution_date = _format_date_range(start_date, end_date)
    
    # Generate unique filename and issue ID
    file_id = str(uuid.uuid4())[:4].upper()
    
    # Generate prefix based on program_title (for ID)
    title_prefix = _get_title_prefix(program_title)
    
    # Use current date for ID (MMYY format)
    date_for_id = datetime.now().strftime("%m%y")
    
    issue_id = f"INT-{title_prefix}-{date_for_id}-{file_id}"
    filename = f"{file_id}.pdf"
    pptx_filename = f"{file_id}.pptx"
    
    # Get directories
    output_dir = _get_output_directory()
    os.makedirs(output_dir, exist_ok=True)
    
    pptx_filepath = os.path.join(output_dir, pptx_filename)
    pdf_filepath = os.path.join(output_dir, filename)
    
    # Get template path
    template_path = os.path.join(
        os.path.dirname(__file__), 
        "..", 
        "static", 
        "certificates", 
        "Template of Certificate Intelligo.pptx"
    )
    template_path = os.path.abspath(template_path)
    
    # Load the presentation
    prs = Presentation(template_path)
    slide = prs.slides[0]
    
    # Define all replacements (will preserve font from template)
    replacements = {
        "{{NAMA}}": name,
        "{{JUDUL}}": program_title,
        "{{PELAKSANAAN}}": execution_date,
        "{{id}}": issue_id
    }
    
    # Replace placeholders in all shapes
    for shape in slide.shapes:
        _replace_text_in_shape(shape, replacements)
    
    # Save the modified PPTX
    prs.save(pptx_filepath)
    
    # Convert PPTX to PDF
    _convert_pptx_to_pdf(pptx_filepath, pdf_filepath)
    
    # Remove temporary PPTX file
    if os.path.exists(pptx_filepath):
        os.remove(pptx_filepath)
    
    return filename


def _convert_pptx_to_pdf(pptx_path: str, pdf_path: str) -> bool:
    """
    Convert PPTX to PDF using LibreOffice (Linux/macOS) or PowerPoint COM (Windows).
    """
    import platform
    
    pptx_path = os.path.abspath(pptx_path)
    pdf_path = os.path.abspath(pdf_path)
    output_dir = os.path.dirname(pdf_path)
    pdf_filename = os.path.basename(pdf_path)
    
    is_windows = platform.system() == "Windows"
    
    if is_windows:
        # Method: Use PowerPoint COM (Windows only)
        try:
            import comtypes.client

            abs_pptx = os.path.abspath(pptx_path)
            abs_pdf = os.path.abspath(pdf_path)
            
            powerpoint = comtypes.client.CreateObject("PowerPoint.Application")
            presentation = powerpoint.Presentations.Open(abs_pptx, WithWindow=False)
            presentation.SaveAs(abs_pdf, 32)  # 32 = ppSaveAsPDF
            presentation.Close()
            powerpoint.Quit()
            
            return True
        except ImportError:
            logger.error("comtypes not installed. Run: pip install comtypes")
        except Exception as e:
            logger.error(f"PowerPoint COM conversion failed: {e}")
        
        raise RuntimeError(
            "Cannot convert PPTX to PDF. Please install comtypes: pip install comtypes"
        )
    else:
        # Method: Use LibreOffice (Linux/macOS)
        try:
            # Use impress_pdf_Export for PowerPoint files (not writer_pdf_Export)
            cmd = [
                "soffice",
                "--headless",
                "--norestore",
                "--nofirststartwizard",
                "--convert-to", "pdf:impress_pdf_Export",
                "--outdir", output_dir,
                pptx_path
            ]
            
            result = subprocess.run(
                cmd,
                capture_output=True,
                text=True,
                timeout=60,
                cwd=output_dir
            )
            
            # Check if PDF was created with expected name
            if os.path.exists(pdf_path):
                logger.info("PDF generated successfully")
                return True
        
        except FileNotFoundError:
            logger.error("LibreOffice (soffice) not found in PATH")
        except Exception as e:
            logger.error(f"LibreOffice conversion error: {e}")
        
        raise RuntimeError(
            "Cannot convert PPTX to PDF. Install LibreOffice: https://www.libreoffice.org/download/"
        )
